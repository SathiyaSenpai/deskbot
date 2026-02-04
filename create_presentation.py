#!/usr/bin/env python3
"""
DeskBot AI Companion - PowerPoint Presentation Generator
For International Innovation Showcase 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BG = RgbColor(26, 26, 46)
CYAN = RgbColor(0, 212, 255)
PURPLE = RgbColor(124, 58, 237)
PINK = RgbColor(244, 114, 182)
GREEN = RgbColor(34, 197, 94)
RED = RgbColor(239, 68, 68)
WHITE = RgbColor(255, 255, 255)
LIGHT_GRAY = RgbColor(148, 163, 184)

def add_title_shape(slide, text, top, font_size=44, color=CYAN, bold=True):
    """Add a title text box"""
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.333), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER
    return txBox

def add_text_box(slide, text, left, top, width, height, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    """Add a text box with formatting"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_bullet_points(slide, items, left, top, width, height, font_size=20, color=WHITE):
    """Add bullet point list"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    return txBox

def set_slide_background(slide, color):
    """Set solid background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
set_slide_background(slide1, DARK_BG)

# Robot emoji/icon
add_text_box(slide1, "🤖", Inches(6), Inches(0.8), Inches(1.5), Inches(1.5), 
             font_size=80, align=PP_ALIGN.CENTER)

# Main title
add_title_shape(slide1, "DeskBot AI Companion", Inches(2), font_size=60, color=CYAN, bold=True)

# Subtitle
add_text_box(slide1, "Your Intelligent Desktop Friend Powered by ESP32 & AI", 
             Inches(1), Inches(3), Inches(11.333), Inches(0.8),
             font_size=28, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Event name
add_text_box(slide1, "🌐 International Innovation Showcase 2026", 
             Inches(2.5), Inches(4.2), Inches(8.333), Inches(0.6),
             font_size=24, color=PINK, bold=True, align=PP_ALIGN.CENTER)

add_text_box(slide1, "Bridging Technology & Human Connection", 
             Inches(2.5), Inches(4.8), Inches(8.333), Inches(0.5),
             font_size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Tech stack boxes at bottom
tech_items = ["🔌 ESP32 Hardware", "🧠 Gemini AI", "🗣️ Voice Interaction", "📱 Mobile Control"]
start_x = 1.5
for i, tech in enumerate(tech_items):
    shape = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                     Inches(start_x + i * 2.7), Inches(6), Inches(2.5), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(40, 40, 70)
    shape.line.color.rgb = RgbColor(80, 80, 120)
    
    tf = shape.text_frame
    tf.paragraphs[0].text = tech
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 2: Problem & Solution
# ============================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide2, DARK_BG)

add_title_shape(slide2, "🎯 The Challenge & Our Solution", Inches(0.3), font_size=44, color=CYAN)

# Problem column (left)
problem_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
problem_box.fill.solid()
problem_box.fill.fore_color.rgb = RgbColor(60, 30, 30)
problem_box.line.color.rgb = RED

add_text_box(slide2, "❌ The Problem", Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.6),
             font_size=28, color=RED, bold=True)

problems = [
    "Increasing isolation in digital workspaces",
    "Mental health challenges from remote work",
    "Expensive AI companions ($500-$2000)",
    "Complex setup requiring technical expertise",
    "No affordable emotional support technology",
    "Limited accessibility for developing regions"
]
add_bullet_points(slide2, problems, Inches(0.8), Inches(2.5), Inches(5.5), Inches(4.2),
                  font_size=18, color=WHITE)

# Solution column (right)
solution_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(6.833), Inches(1.5), Inches(6), Inches(5.5))
solution_box.fill.solid()
solution_box.fill.fore_color.rgb = RgbColor(20, 50, 30)
solution_box.line.color.rgb = GREEN

add_text_box(slide2, "✅ DeskBot Solution", Inches(7.133), Inches(1.7), Inches(5.5), Inches(0.6),
             font_size=28, color=GREEN, bold=True)

solutions = [
    "Affordable companion under $30 total cost",
    "Emotional expressions via OLED & LED lights",
    "Free AI powered by Google Gemini API",
    "Plug-and-play with phone hotspot",
    "Multilingual support (English + Tamil)",
    "DIY-friendly cardboard/3D printed body"
]
add_bullet_points(slide2, solutions, Inches(7.133), Inches(2.5), Inches(5.5), Inches(4.2),
                  font_size=18, color=WHITE)

# ============================================================================
# SLIDE 3: Key Features
# ============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide3, DARK_BG)

add_title_shape(slide3, "✨ Key Features & Capabilities", Inches(0.2), font_size=44, color=CYAN)

features = [
    ("👀", "Expressive OLED Eyes", "Custom eye engine with 15+ emotional states"),
    ("🧠", "AI Conversation", "Google Gemini 2.5 Flash for intelligent chat"),
    ("🗣️", "Voice Interaction", "Edge TTS output + Web Speech input"),
    ("🌈", "Mood LED Ring", "16 WS2812 LEDs synced with emotions"),
    ("📱", "Mobile Web Control", "Responsive dashboard from any browser"),
    ("🤲", "Multi-Sensor Aware", "Touch, motion, distance, light sensors"),
]

# Create 2x3 grid of feature boxes
for i, (icon, title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(1.4 + row * 3)
    
    # Feature box
    box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.9), Inches(2.7))
    box.fill.solid()
    box.fill.fore_color.rgb = RgbColor(35, 35, 55)
    box.line.color.rgb = RgbColor(60, 60, 90)
    
    # Icon
    add_text_box(slide3, icon, left + Inches(1.4), top + Inches(0.2), Inches(1), Inches(0.8),
                 font_size=48, align=PP_ALIGN.CENTER)
    
    # Title
    add_text_box(slide3, title, left + Inches(0.2), top + Inches(1.1), Inches(3.5), Inches(0.5),
                 font_size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    # Description
    add_text_box(slide3, desc, left + Inches(0.2), top + Inches(1.7), Inches(3.5), Inches(0.8),
                 font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 4: Technical Architecture
# ============================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide4, DARK_BG)

add_title_shape(slide4, "🏗️ System Architecture", Inches(0.2), font_size=44, color=CYAN)

# Three architecture boxes
arch_data = [
    (PURPLE, "🔌 ESP32 Hardware", ["SH1106 OLED Display", "SG90 Servo Motor", "WS2812 LED Ring", "Ultrasonic Sensor", "PIR Motion Sensor", "Touch Sensors"]),
    (RgbColor(8, 145, 178), "🖥️ Node.js Server", ["WebSocket Bridge", "Gemini AI Integration", "Edge TTS Engine", "Audio Management", "State Synchronization"]),
    (RgbColor(5, 150, 105), "🌐 Web Interface", ["Real-time Dashboard", "Chat Interface", "Behavior Controls", "Sensor Monitoring", "Speech Recognition"]),
]

for i, (color, title, items) in enumerate(arch_data):
    left = Inches(0.7 + i * 4.3)
    
    # Box
    box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(3.5), Inches(4))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = WHITE
    
    # Title
    add_text_box(slide4, title, left + Inches(0.1), Inches(1.7), Inches(3.3), Inches(0.6),
                 font_size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    # Items
    add_bullet_points(slide4, items, left + Inches(0.2), Inches(2.4), Inches(3.1), Inches(3),
                      font_size=14, color=WHITE)

# Arrows between boxes
add_text_box(slide4, "⟷", Inches(4), Inches(3.2), Inches(0.6), Inches(0.6),
             font_size=36, color=CYAN, align=PP_ALIGN.CENTER)
add_text_box(slide4, "⟷", Inches(8.3), Inches(3.2), Inches(0.6), Inches(0.6),
             font_size=36, color=CYAN, align=PP_ALIGN.CENTER)

# Tech stack at bottom
tech_stack = ["⚡ C++ / Arduino", "🟢 Node.js", "🌐 WebSocket", "🤖 Gemini 2.5", "🔊 Edge TTS", "📱 PWA Ready"]
for i, tech in enumerate(tech_stack):
    shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.5 + i * 2.1), Inches(6.2), Inches(2), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(40, 40, 60)
    shape.line.color.rgb = RgbColor(80, 80, 100)
    
    tf = shape.text_frame
    tf.paragraphs[0].text = tech
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 5: Impact & Future
# ============================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide5, DARK_BG)

add_title_shape(slide5, "🌍 Impact & Future Vision", Inches(0.2), font_size=44, color=CYAN)

# Social Impact section (left)
add_text_box(slide5, "🎯 Social Impact", Inches(0.5), Inches(1.2), Inches(6), Inches(0.5),
             font_size=24, color=GREEN, bold=True)

impacts = [
    "💰 Democratizes AI for all income levels",
    "🧠 Mental wellness through friendly interaction",
    "🎓 Educational tool for STEM & robotics",
    "👴 Elderly companionship & assistance",
    "🌐 Multilingual accessibility"
]
add_bullet_points(slide5, impacts, Inches(0.5), Inches(1.8), Inches(6), Inches(2.5),
                  font_size=16, color=WHITE)

# Future Roadmap section (right)
add_text_box(slide5, "🚀 Future Roadmap", Inches(6.833), Inches(1.2), Inches(6), Inches(0.5),
             font_size=24, color=PINK, bold=True)

future = [
    "🏥 Healthcare integration for patients",
    "📚 Educational AI tutor capabilities",
    "🏠 Smart home IoT hub integration",
    "🎮 Gamification for productivity",
    "🤝 Open-source community expansion"
]
add_bullet_points(slide5, future, Inches(6.833), Inches(1.8), Inches(6), Inches(2.5),
                  font_size=16, color=WHITE)

# CTA Box at bottom
cta_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(1.5), Inches(4.8), Inches(10.333), Inches(2.2))
cta_box.fill.solid()
cta_box.fill.fore_color.rgb = RgbColor(30, 40, 60)
cta_box.line.color.rgb = CYAN

add_text_box(slide5, "🤖 Meet DeskBot - Technology with Heart", 
             Inches(1.5), Inches(5), Inches(10.333), Inches(0.6),
             font_size=32, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

add_text_box(slide5, "Affordable • Open Source • Accessible • Human-Centered AI", 
             Inches(1.5), Inches(5.7), Inches(10.333), Inches(0.5),
             font_size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text_box(slide5, "✨ Thank you! Questions Welcome! ✨", 
             Inches(1.5), Inches(6.3), Inches(10.333), Inches(0.5),
             font_size=24, color=PINK, bold=True, align=PP_ALIGN.CENTER)

# ============================================================================
# Save the presentation
# ============================================================================
output_path = "/home/sathiya/Documents/Project_Companion/deskbot/DeskBot_International_Presentation.pptx"
prs.save(output_path)
print(f"✅ PowerPoint presentation saved to: {output_path}")
